from __future__ import annotations

import hashlib
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation


@dataclass
class SlideDoc:
    slide_id: str
    text: str
    metadata: Dict[str, Any]


def _shape_text(shape) -> str:
    try:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            return shape.text or ""
    except Exception:
        pass
    return ""


def _notes_text(slide) -> str:
    try:
        notes = slide.notes_slide
        if notes and notes.notes_text_frame:
            return notes.notes_text_frame.text or ""
    except Exception:
        pass
    return ""


def _text_hash(text: str) -> str:
    return hashlib.sha1(text.encode("utf-8")).hexdigest()


def extract_slides(
    pptx_path: str | Path,
    doc_hash: str,
    source_root: str | None = None,
    doc_uid: str | None = None,
    doc_signature: str | None = None,
) -> List[SlideDoc]:
    p = Path(pptx_path).expanduser().resolve()
    prs = Presentation(str(p))

    deck_title = p.stem
    num_slides = len(prs.slides)
    slide_docs: List[SlideDoc] = []

    for idx, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            text = _shape_text(shape).strip()
            if text:
                parts.append(text)

        notes = _notes_text(slide).strip()
        if notes:
            parts.append("NOTES:\n" + notes)

        full_text = "\n".join(parts).strip()
        if not full_text:
            continue

        stable_doc_uid = doc_uid or hashlib.sha1(str(p).encode("utf-8")).hexdigest()
        slide_id = f"{stable_doc_uid}::slide::{idx}"
        metadata: Dict[str, Any] = {
            "doc_hash": doc_hash,
            "doc_uid": stable_doc_uid,
            "doc_signature": doc_signature or "",
            "source_root": source_root if (source_root and source_root.startswith(("sharepoint::", "local::"))) else (str(Path(source_root).expanduser().resolve()) if source_root else str(p.parent)),
            "deck_title": deck_title,
            "pptx_path": str(p),
            "slide_number": idx,
            "num_slides": num_slides,
            "slide_text_hash": _text_hash(full_text),
        }
        slide_docs.append(SlideDoc(slide_id=slide_id, text=full_text, metadata=metadata))

    return slide_docs
