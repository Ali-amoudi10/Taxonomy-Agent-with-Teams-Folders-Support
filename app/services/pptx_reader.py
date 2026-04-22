from __future__ import annotations
import json
from pathlib import Path
from typing import Dict, Any

from pptx import Presentation

_CACHE: Dict[str, Dict[str, Any]] = {}
_CACHE_LOADED_FROM: str | None = None

def _load_cache(cache_path: str) -> None:
    global _CACHE_LOADED_FROM, _CACHE
    if _CACHE_LOADED_FROM == cache_path:
        return
    _CACHE_LOADED_FROM = cache_path
    p = Path(cache_path)
    p.parent.mkdir(parents=True, exist_ok=True)
    if not p.exists():
        _CACHE = {}
        return
    try:
        _CACHE = json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        _CACHE = {}

def _save_cache(cache_path: str) -> None:
    p = Path(cache_path)
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(_CACHE, ensure_ascii=False), encoding="utf-8")

def extract_text(pptx_path: str, cache_path: str, max_chars: int = 40000) -> str:
    """
    Extracts text from slides + speaker notes (best-effort).
    Uses a simple JSON cache keyed by absolute path and mtime.
    """
    _load_cache(cache_path)

    path = Path(pptx_path).resolve()
    key = str(path)
    try:
        mtime = path.stat().st_mtime
    except FileNotFoundError:
        return ""

    cached = _CACHE.get(key)
    if cached and cached.get("mtime") == mtime and isinstance(cached.get("text"), str):
        return cached["text"][:max_chars]

    text_parts: list[str] = []
    try:
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, start=1):
            text_parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    t = shape.text.strip()
                    if t:
                        text_parts.append(t)

            # notes
            try:
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    nt = (slide.notes_slide.notes_text_frame.text or "").strip()
                    if nt:
                        text_parts.append(f"[Notes {i}] {nt}")
            except Exception:
                pass
    except Exception:
        # corrupted/unreadable files should not crash the whole scan
        extracted = ""
    else:
        extracted = "\n".join(text_parts)

    _CACHE[key] = {"mtime": mtime, "text": extracted}
    try:
        _save_cache(cache_path)
    except Exception:
        pass

    return extracted[:max_chars]