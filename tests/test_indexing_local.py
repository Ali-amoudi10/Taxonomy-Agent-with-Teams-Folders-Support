from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.rag_kb.chroma_store import ChromaSlideStore
from app.rag_kb.config import KBConfig
from app.rag_kb.indexer import index_root
from app.rag_kb.retriever import search_slides


def _fake_embed(self, texts):
    out = []
    for text in texts:
        text = (text or "").lower()
        out.append([
            1.0 if "finance" in text else 0.0,
            1.0 if "healthcare" in text else 0.0,
            float(len(text) % 13) / 13.0,
        ])
    return out


def _make_pptx(path: Path, slides: list[str]):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for idx, body in enumerate(slides, start=1):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {idx}"
        slide.placeholders[1].text = body
    prs.save(path)


def test_local_index_and_search(monkeypatch, tmp_path):
    monkeypatch.setattr(ChromaSlideStore, "embed_texts", _fake_embed, raising=False)
    source_dir = tmp_path / "slides"
    source_dir.mkdir()
    _make_pptx(source_dir / "finance_deck.pptx", ["finance budgeting strategy", "healthcare operations"])
    _make_pptx(source_dir / "general_deck.pptx", ["talent management", "delivery excellence"])

    cfg = KBConfig(chroma_dir=str(tmp_path / "chroma"))
    result = index_root(str(source_dir), cfg)
    assert result["ok"] is True
    assert result["stats"]["indexed_files"] == 2

    hits = search_slides("finance", cfg, top_k=5)
    assert hits
    assert hits[0]["deck_title"] == "finance_deck"
